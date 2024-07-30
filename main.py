import streamlit as st
import base64
from openai import OpenAI
import os
import tempfile
import io
from PyPDF2 import PdfReader
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from dotenv import load_dotenv

load_dotenv()

api_key = os.getenv("OPENAI_API_KEY")
print(f"API Key: {api_key}")  # This should print your API key or 'None' if it's not set

if not api_key:
    raise ValueError("OPENAI_API_KEY environment variable is not set")

client = OpenAI(api_key=api_key)

CHUNK_SIZE = 20000  # Adjust this value as needed to fit within the token limit

def extract_text_from_pdf(pdf_file):
    pdf_reader = PdfReader(pdf_file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text() + "\n"
    return text

def create_assistant():
    assistant = client.beta.assistants.create(
        instructions="You are a presentation creation assistant. When given content from a PDF, create a PowerPoint presentation summarizing the key points.",
        model="gpt-4o-mini",
        tools=[{"type": "code_interpreter"}]
    )
    print(f"Assistant created: {assistant}")
    return assistant

def create_thread_and_run(assistant_id, content):
    # Create the thread
    thread = client.beta.threads.create()
    print(f"Thread created: {thread}")

    # Submit the message with detailed instructions
    client.beta.threads.messages.create(
        thread_id=thread.id, role="user", content=f"""Create a PowerPoint presentation summarizing the following content:\n\n{content}\n\nFor the presentation:\n1. Create a title slide with an appropriate title and subtitle based on the content.\n2. Create maximum 50 content slides summarizing key points.\n3. Use appropriate formatting and structure to make the presentation visually appealing.\n4. For each slide, provide:\n   - The slide title\n   - The slide content (bullet points or short paragraphs)\n\nOutput the slide content in a structured format that can be easily parsed, like this:\n\nSLIDE 1 (Title Slide):\nTitle: [Title]\nSubtitle: [Subtitle]\n\nSLIDE 2:\nTitle: [Slide Title]\nContent:\n- [Bullet point 1]\n- [Bullet point 2]\n- [Bullet point 3]\n\n[Continue for all slides]"""
    )

    # Create the run
    run = client.beta.threads.runs.create(
        thread_id=thread.id,
        assistant_id=assistant_id
    )
    print(f"Run created: {run}")
    return thread, run

def wait_for_run_completion(thread_id, run_id, timeout=1200, check_interval=10):
    start_time = time.time()
    while time.time() - start_time < timeout:
        run = client.beta.threads.runs.retrieve(thread_id=thread_id, run_id=run_id)
        print(f"Run status: {run.status}")
        if run.status == "succeeded":
            return run
        elif run.status == "failed":
            print(f"Run failed: {run}")
            raise Exception("Run failed")
        elif run.status == "completed":
            return run
        time.sleep(check_interval)
    raise TimeoutError("Run did not complete in time")

def get_presentation_content(thread_id):
    thread = client.beta.threads.retrieve(thread_id)
    print(f"Thread ID: {thread.id}")
    print(f"Tool Resources: {thread.tool_resources}")

    # Check if there are messages in the thread
    messages = client.beta.threads.messages.list(thread_id=thread.id)
    print(f"Messages: {messages}")

    for message in messages.data:
        if message.role == 'assistant':
            for content_block in message.content:
                if content_block.type == 'text':
                    return content_block.text.value

    raise AttributeError("Thread object has no attribute 'messages'")

def create_presentation(content):
    print("Creating presentation from content...")
    print(content)  # Debug print to show the content structure

    slides = content.split("---\n\n")  # Adjust the split based on the separator used in the assistant's response
    print(f"Total slides parsed: {len(slides)}")

    prs = Presentation()

    for slide_content in slides:
        if not slide_content.strip():
            continue

        print(f"Processing slide content: {slide_content}")  # Debug print to show the slide content being processed

        # Split the slide content into lines
        lines = slide_content.split("\n")
        title = ""
        subtitle = ""
        content = []

        for line in lines:
            if line.startswith("Title: "):
                title = line[len("Title: "):]
            elif line.startswith("Subtitle: "):
                subtitle = line[len("Subtitle: "):]
            elif line.startswith("Content:"):
                content.append(line[len("Content: "):])
            elif line.startswith("- "):
                content.append(line)

        # Add slide to presentation
        slide_layout = prs.slide_layouts[1] if "Title Slide" not in slide_content else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        title_placeholder = slide.shapes.title
        title_placeholder.text = title

        if slide_layout == prs.slide_layouts[0]:  # Title Slide
            subtitle_placeholder = slide.placeholders[1]
            subtitle_placeholder.text = subtitle
        else:
            content_placeholder = slide.placeholders[1].text_frame
            for line in content:
                p = content_placeholder.add_paragraph()
                p.text = line
                p.font.size = Pt(18)

            # Apply some formatting to make it visually appealing
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.name = 'Calibri'
                    paragraph.font.size = Pt(18)
                    paragraph.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

            # Center align the content
            for paragraph in content_placeholder.paragraphs:
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
        prs.save(tmp_file.name)
        return tmp_file.name

def main():
    st.title("PowerPoint Presentation Generator from PDF")

    uploaded_file = st.file_uploader("Upload a PDF file", type="pdf")
    generate_button = st.button("Generate Presentation")

    if generate_button and uploaded_file:
        st.info("Generating presentation... Please wait.")

        # Extract text from the PDF
        text_content = extract_text_from_pdf(uploaded_file)

        # Split the text content into chunks
        text_chunks = [text_content[i:i + CHUNK_SIZE] for i in range(0, len(text_content), CHUNK_SIZE)]

        # Create assistant
        assistant = create_assistant()

        all_presentation_content = []

        for chunk in text_chunks:
            # Create thread and run for each chunk
            thread, run = create_thread_and_run(assistant.id, chunk)

            # Wait for run completion
            try:
                run = wait_for_run_completion(thread.id, run.id)
            except Exception as e:
                st.error(f"Failed to complete the run for a chunk: {e}")
                return

            # Get the presentation content for the chunk
            try:
                presentation_content = get_presentation_content(thread.id)
                all_presentation_content.append(presentation_content)
            except AttributeError as e:
                st.error(f"Error retrieving presentation content for a chunk: {e}")
                return

        combined_content = "\n\n---\n\n".join(all_presentation_content)

        if combined_content:
            st.success("Presentation content generated successfully!")
            
            # Create the presentation
            try:
                ppt_file_path = create_presentation(combined_content)
            except Exception as e:
                st.error(f"Error creating presentation: {e}")
                return
            
            # Provide download link
            st.markdown(get_ppt_download_link(ppt_file_path), unsafe_allow_html=True)
            
            # Clean up
            os.unlink(ppt_file_path)
        else:
            st.error("Failed to generate presentation. Please try again.")

        # Clean up
        client.beta.assistants.delete(assistant.id)

def get_ppt_download_link(file_path):
    with open(file_path, "rb") as file:
        ppt_contents = file.read()
    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="presentation.pptx">Download the PowerPoint Presentation</a>'

if __name__ == "__main__":
    main()
